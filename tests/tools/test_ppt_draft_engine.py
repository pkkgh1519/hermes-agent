import base64
from pathlib import Path
from xml.sax.saxutils import escape
from zipfile import ZipFile

import pytest
from PIL import Image

from gateway.ppt_draft_state import DraftPhotoBatch
from tools import ppt_draft_engine as draft_engine
from tools.ppt_draft_engine import (
    DraftInputError,
    build_draft_payload,
    create_draft_pptx,
    parse_offers_csv,
)


REQUIRED_HEADER = "id,name,location,size_floor,price,points,photo_tag\n"
PNG_1X1 = base64.b64decode(
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+aF9sAAAAASUVORK5CYII="
)


def _write_csv(path: Path, content: str) -> Path:
    path.write_text(content, encoding="utf-8")
    return path


def _xlsx_column_name(index: int) -> str:
    name = ""
    while index > 0:
        index, remainder = divmod(index - 1, 26)
        name = chr(65 + remainder) + name
    return name


def _write_xlsx(path: Path, rows: list[list[str]], *, sheet_name: str = "in") -> Path:
    worksheet_rows: list[str] = []
    for row_idx, row in enumerate(rows, start=1):
        cells: list[str] = []
        for col_idx, value in enumerate(row, start=1):
            cell_ref = f"{_xlsx_column_name(col_idx)}{row_idx}"
            cells.append(
                f'<c r="{cell_ref}" t="inlineStr"><is><t xml:space="preserve">{escape(str(value or ""))}</t></is></c>'
            )
        worksheet_rows.append(f'<row r="{row_idx}">{"".join(cells)}</row>')

    workbook_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<workbook xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main" '
        'xmlns:r="http://schemas.openxmlformats.org/officeDocument/2006/relationships">'
        '<sheets>'
        f'<sheet name="{escape(sheet_name)}" sheetId="1" r:id="rId1"/>'
        '</sheets>'
        '</workbook>'
    )
    worksheet_xml = (
        '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
        '<worksheet xmlns="http://schemas.openxmlformats.org/spreadsheetml/2006/main">'
        f'<sheetData>{"".join(worksheet_rows)}</sheetData>'
        '</worksheet>'
    )

    with ZipFile(path, "w") as zf:
        zf.writestr(
            "[Content_Types].xml",
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">'
            '<Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>'
            '<Default Extension="xml" ContentType="application/xml"/>'
            '<Override PartName="/xl/workbook.xml" '
            'ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet.main+xml"/>'
            '<Override PartName="/xl/worksheets/sheet1.xml" '
            'ContentType="application/vnd.openxmlformats-officedocument.spreadsheetml.worksheet+xml"/>'
            '</Types>',
        )
        zf.writestr(
            "_rels/.rels",
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" '
            'Target="xl/workbook.xml"/>'
            '</Relationships>',
        )
        zf.writestr("xl/workbook.xml", workbook_xml)
        zf.writestr(
            "xl/_rels/workbook.xml.rels",
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">'
            '<Relationship Id="rId1" '
            'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/worksheet" '
            'Target="worksheets/sheet1.xml"/>'
            '</Relationships>',
        )
        zf.writestr("xl/worksheets/sheet1.xml", worksheet_xml)
    return path


def _write_image(path: Path, *, size: tuple[int, int] = (1, 1), color: tuple[int, int, int] = (240, 240, 240)) -> str:
    if size == (1, 1):
        path.write_bytes(PNG_1X1)
    else:
        Image.new("RGB", size, color).save(path)
    return str(path)


def _offer_csv_row(
    offer_id: str,
    *,
    name: str | None = None,
    location: str | None = None,
    size_floor: str = "2층 / 12평",
    price: str = "보증금 2200 / 월세 220 / 관리비 35",
    points: str = "성수역 도보 6분; 무료주차 1대 가능; 지상 로딩도크 및 화물 E/V 有",
    photo_tag: str | None = None,
) -> str:
    return ",".join(
        [
            offer_id,
            name or f"건물 {offer_id}",
            location or f"성수동 {offer_id}",
            size_floor,
            price,
            points,
            photo_tag or offer_id,
        ]
    )


def _count_slides(pptx_path: Path) -> int:
    with ZipFile(pptx_path) as zf:
        return len([name for name in zf.namelist() if name.startswith("ppt/slides/slide") and name.endswith(".xml")])


def test_parse_offers_csv_rejects_missing_required_columns(tmp_path):
    csv_path = _write_csv(
        tmp_path / "offers.csv",
        "id,name,location,size_floor,price,points\n"
        "offer_01,캐슬프라자 3층,성수동,3층 / 18평,2000/160/42,주차 1대\n",
    )

    with pytest.raises(DraftInputError) as exc:
        parse_offers_csv(csv_path)

    assert exc.value.code == "missing_columns"
    assert "photo_tag" in str(exc.value)



def test_parse_offers_csv_normalizes_whitespace_padded_headers(tmp_path):
    csv_path = _write_csv(
        tmp_path / "offers.csv",
        "id , name , location , size_floor , price , points , photo_tag \n"
        "offer_01,서울숲드림타워,성수동2가323,2층 / 12평,보증금 2200 / 월세 220 / 관리비 35,성수역 도보 6분; 무료주차 1대 가능,offer_01\n",
    )

    offers = parse_offers_csv(csv_path)

    assert [offer.id for offer in offers] == ["offer_01"]
    assert offers[0].photo_tag == "offer_01"
    assert offers[0].name == "서울숲드림타워"



def test_parse_offers_csv_accepts_xlsx_and_splits_newline_points(tmp_path):
    xlsx_path = _write_xlsx(
        tmp_path / "offers.xlsx",
        [
            ["id", "name", "location", "size_floor", "price", "points", "photo_tag"],
            [
                "offer_01",
                "서울숲드림타워",
                "성수동2가323",
                "2층 / 12평",
                "보증금 2200 / 월세 220 / 관리비 35",
                "성수역 도보 6분\n무료주차 1대 가능\n지상 로딩도크 및 화물 E/V 有",
                "offer_01",
            ],
        ],
    )

    offers = parse_offers_csv(xlsx_path)

    assert [offer.id for offer in offers] == ["offer_01"]
    assert offers[0].points == [
        "성수역 도보 6분",
        "무료주차 1대 가능",
        "지상 로딩도크 및 화물 E/V 有",
    ]
    assert offers[0].photo_tag == "offer_01"



def test_parse_offers_csv_wraps_malformed_xlsx_xml_in_draft_input_error(tmp_path):
    xlsx_path = tmp_path / "broken.xlsx"
    with ZipFile(xlsx_path, "w") as zf:
        zf.writestr("xl/workbook.xml", "<broken")
        zf.writestr(
            "xl/_rels/workbook.xml.rels",
            '<?xml version="1.0" encoding="UTF-8" standalone="yes"?>'
            '<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships"></Relationships>',
        )

    with pytest.raises(DraftInputError) as exc:
        parse_offers_csv(xlsx_path)

    assert exc.value.code == "invalid_xlsx"



def test_build_draft_payload_rejects_unmatched_photo_tag(tmp_path):
    csv_path = _write_csv(
        tmp_path / "offers.csv",
        REQUIRED_HEADER
        + "offer_01,캐슬프라자 3층,성수동 2가 289-13 / 성수역 도보 2분,3층 / 18평,보증금 2000 / 월세 160 / 관리비 42,주차 1대; 공실,offer_01\n",
    )

    with pytest.raises(DraftInputError) as exc:
        build_draft_payload(
            csv_path,
            photo_batches=[
                DraftPhotoBatch(
                    tag="offer_02",
                    image_paths=["/tmp/offer_02_01.jpg"],
                    message_id="m1",
                    uploaded_at=1710000000.0,
                )
            ],
        )

    assert exc.value.code == "photo_tag_not_matched"
    assert "offer_01" in str(exc.value)



def test_build_draft_payload_normalizes_points_and_aggregates_photo_batches(tmp_path):
    csv_path = _write_csv(
        tmp_path / "offers.csv",
        REQUIRED_HEADER
        + "offer_01,캐슬프라자 3층,성수동 2가 289-13 / 성수역 도보 2분,3층 / 18평,보증금 2000 / 월세 160 / 관리비 42,주차 1대; 인테리어 있음; 공실,offer_01\n"
        + "offer_02,동성빌딩 4층,성수동 2가 277-50 / 성수역 도보 4분,4층 / 14평,보증금 1100 / 월세 105 / 관리비 16,주차 가능; 공실; 화물 EV 가능,offer_02\n",
    )

    payload = build_draft_payload(
        csv_path,
        photo_batches=[
            DraftPhotoBatch(
                tag=" offer_01 ",
                image_paths=["/tmp/offer_01_01.jpg", "/tmp/offer_01_02.jpg"],
                message_id="m1",
                uploaded_at=1710000000.0,
            ),
            DraftPhotoBatch(
                tag="offer_01",
                image_paths=["/tmp/offer_01_03.jpg"],
                message_id="m2",
                uploaded_at=1710000001.0,
            ),
            DraftPhotoBatch(
                tag="offer_02",
                image_paths=["/tmp/offer_02_01.jpg"],
                message_id="m3",
                uploaded_at=1710000002.0,
            ),
        ],
    )

    assert payload.csv_path == str(csv_path)
    assert payload.matched_photo_count == 4
    assert payload.unmatched_photo_tags == []
    assert [offer.id for offer in payload.offers] == ["offer_01", "offer_02"]
    assert payload.offers[0].points == ["주차 1대", "인테리어 있음", "공실"]
    assert payload.offers[0].photo_paths == [
        "/tmp/offer_01_01.jpg",
        "/tmp/offer_01_02.jpg",
        "/tmp/offer_01_03.jpg",
    ]
    assert payload.offers[1].photo_paths == ["/tmp/offer_02_01.jpg"]



def test_draft_template_config_defaults():
    config = draft_engine.DraftTemplateConfig()

    assert config.briefing_offers_per_slide == 5
    assert config.detail_images_per_slide == 2
    assert config.gallery_images_per_slide == 9



def test_draft_template_config_normalizes_intro_assets_to_paths(tmp_path):
    config = draft_engine.DraftTemplateConfig(
        intro_assets=(str(tmp_path / "intro-1.png"), str(tmp_path / "intro-2.png")),
    )

    assert all(isinstance(path, Path) for path in config.intro_assets)



def test_draft_template_config_rejects_unsupported_values():
    with pytest.raises(ValueError):
        draft_engine.DraftTemplateConfig(briefing_offers_per_slide=0)

    with pytest.raises(ValueError):
        draft_engine.DraftTemplateConfig(gallery_images_per_slide=0)

    with pytest.raises(ValueError):
        draft_engine.DraftTemplateConfig(detail_images_per_slide=3)



def test_build_deck_render_plan_chunks_briefing_by_five(tmp_path):
    rows = [REQUIRED_HEADER.strip()]
    for idx in range(1, 7):
        offer_id = f"offer_{idx:02d}"
        rows.append(_offer_csv_row(offer_id))
    csv_path = _write_csv(tmp_path / "offers.csv", "\n".join(rows) + "\n")

    payload = build_draft_payload(
        csv_path,
        photo_batches=[
            DraftPhotoBatch(tag=f"offer_{idx:02d}", image_paths=[f"/tmp/offer_{idx:02d}.jpg"])
            for idx in range(1, 7)
        ],
    )

    plan = draft_engine.build_deck_render_plan(payload, draft_engine.DraftTemplateConfig())

    assert len(plan.briefing_chunks) == 2
    assert [len(chunk) for chunk in plan.briefing_chunks] == [5, 1]



def test_build_offer_render_plan_prefers_unique_detail_images_before_reusing_overview(tmp_path):
    csv_path = _write_csv(
        tmp_path / "offers.csv",
        REQUIRED_HEADER + _offer_csv_row("offer_01") + "\n",
    )
    image_dir = tmp_path / "images"
    image_dir.mkdir()

    hero_image = _write_image(image_dir / "hero.png", size=(1600, 900), color=(40, 70, 120))
    detail_one = _write_image(image_dir / "detail-01.png", size=(1400, 900), color=(120, 130, 140))
    detail_two = _write_image(image_dir / "detail-02.png", size=(1500, 900), color=(160, 170, 180))
    map_image = _write_image(image_dir / "map.png", size=(900, 1400), color=(230, 230, 230))
    gallery_image = _write_image(image_dir / "gallery.png", size=(1500, 900), color=(200, 210, 220))

    payload = build_draft_payload(
        csv_path,
        photo_batches=[
            DraftPhotoBatch(
                tag="offer_01",
                image_paths=[hero_image, detail_one, detail_two, map_image, gallery_image],
            )
        ],
    )

    offer = payload.offers[0]
    plan = draft_engine.build_offer_render_plan(offer, draft_engine.DraftTemplateConfig())
    gallery_paths = [path for chunk in plan.gallery_chunks for path in chunk]

    assert plan.hero_image == hero_image
    assert plan.map_image == map_image
    assert plan.detail_images == [detail_one, detail_two]
    assert gallery_paths == [gallery_image]
    assert plan.hero_image not in plan.detail_images
    assert plan.map_image not in plan.detail_images



def test_build_offer_overview_layout_wraps_long_headline_and_separates_rows():
    offer = draft_engine.ParsedOffer(
        id="offer_01",
        name="동성빌딩 4층 임대 오피스",
        location="성수동 2가 277-50 / 성수역 도보 4분 / 서울숲 생활권",
        size_floor="4층 / 14평",
        price="보증금 200 / 월세 105 / 관리비 16",
        points=["주차 가능", "공실", "화물 EV 가능"],
        photo_tag="offer_01",
    )

    layout = draft_engine.build_offer_overview_layout(offer)

    assert "\n" in layout.headline_text
    assert layout.price_row_top > layout.subtitle_top
    assert layout.note_row_top > layout.price_row_top
    assert layout.photos_top > layout.note_row_top



def test_create_draft_pptx_uses_safe_overview_layout_for_multiple_long_title_offers(tmp_path):
    from pptx import Presentation

    csv_path = _write_csv(
        tmp_path / "offers.csv",
        REQUIRED_HEADER
        + _offer_csv_row(
            "offer_01",
            name="동성빌딩 4층 임대 오피스",
            location="성수동 2가 277-50 / 성수역 도보 4분 / 서울숲 생활권",
        )
        + "\n"
        + _offer_csv_row(
            "offer_02",
            name="서울숲프라자 6층 업무시설",
            location="성수동 2가 288-11 / 성수역 도보 3분 / 메인대로 코너 입지",
            price="보증금 350 / 월세 180 / 관리비 22",
            points="주차 가능; 공실; 천장형 냉난방기",
        )
        + "\n",
    )
    image_dir = tmp_path / "images"
    image_dir.mkdir()
    payload = build_draft_payload(
        csv_path,
        photo_batches=[
            DraftPhotoBatch(
                tag="offer_01",
                image_paths=[
                    _write_image(image_dir / "offer_01_01.png", size=(1600, 900), color=(40, 70, 120)),
                    _write_image(image_dir / "offer_01_02.png", size=(1200, 800), color=(180, 190, 205)),
                    _write_image(image_dir / "offer_01_03.png", size=(1500, 900), color=(230, 230, 230)),
                ],
            ),
            DraftPhotoBatch(
                tag="offer_02",
                image_paths=[
                    _write_image(image_dir / "offer_02_01.png", size=(1600, 900), color=(70, 90, 120)),
                    _write_image(image_dir / "offer_02_02.png", size=(1200, 800), color=(190, 200, 210)),
                    _write_image(image_dir / "offer_02_03.png", size=(1500, 900), color=(210, 220, 230)),
                ],
            ),
        ],
    )

    output_path = tmp_path / "draft-multi-long-overview.pptx"
    create_draft_pptx(payload, output_path)

    prs = Presentation(str(output_path))
    overview_slides = []
    for slide in prs.slides:
        slide_texts = [
            paragraph.text
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            for paragraph in shape.text_frame.paragraphs
            if paragraph.text
        ]
        if any(text.startswith("02. 매물 소개") for text in slide_texts) and "보증금" in slide_texts:
            overview_slides.append(slide)

    assert len(overview_slides) == 2

    observed_price_tops = []
    for slide in overview_slides:
        price_shape = next(
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and any(paragraph.text == "보증금" for paragraph in shape.text_frame.paragraphs)
        )
        note_shape = next(
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and any(paragraph.text == "비고" for paragraph in shape.text_frame.paragraphs)
        )
        headline_shape = next(
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and any(
                ("동성빌딩" in paragraph.text) or ("서울숲프라자" in paragraph.text)
                for paragraph in shape.text_frame.paragraphs
            )
        )
        observed_price_tops.append(price_shape.top)
        assert price_shape.top >= headline_shape.top + headline_shape.height
        assert note_shape.top > price_shape.top

    assert len(set(observed_price_tops)) == 1



def test_create_draft_pptx_accepts_custom_template_config(tmp_path):
    rows = [REQUIRED_HEADER.strip()]
    image_dir = tmp_path / "images"
    image_dir.mkdir()
    photo_batches: list[DraftPhotoBatch] = []
    for idx in range(1, 7):
        offer_id = f"offer_{idx:02d}"
        rows.append(_offer_csv_row(offer_id))
        photo_batches.append(
            DraftPhotoBatch(
                tag=offer_id,
                image_paths=[_write_image(image_dir / f"{offer_id}.png")],
            )
        )
    csv_path = _write_csv(tmp_path / "offers.csv", "\n".join(rows) + "\n")
    payload = build_draft_payload(csv_path, photo_batches=photo_batches)

    output_path = tmp_path / "draft-custom-config.pptx"
    create_draft_pptx(
        payload,
        output_path,
        template_config=draft_engine.DraftTemplateConfig(briefing_offers_per_slide=10),
    )

    assert output_path.exists()
    assert _count_slides(output_path) == 16



def test_create_draft_pptx_rejects_non_image_intro_assets(tmp_path):
    csv_path = _write_csv(
        tmp_path / "offers.csv",
        REQUIRED_HEADER + _offer_csv_row("offer_01") + "\n",
    )
    image_dir = tmp_path / "images"
    image_dir.mkdir()
    payload = build_draft_payload(
        csv_path,
        photo_batches=[
            DraftPhotoBatch(
                tag="offer_01",
                image_paths=[_write_image(image_dir / "offer_01.png")],
            )
        ],
    )

    intro_one = tmp_path / "intro-1.txt"
    intro_two = tmp_path / "intro-2.txt"
    intro_one.write_text("not an image", encoding="utf-8")
    intro_two.write_text("also not an image", encoding="utf-8")

    with pytest.raises(DraftInputError) as exc:
        create_draft_pptx(
            payload,
            tmp_path / "draft-invalid-intro.pptx",
            template_config=draft_engine.DraftTemplateConfig(
                intro_assets=(intro_one, intro_two),
            ),
        )

    assert exc.value.code == "invalid_intro_assets"



def test_create_draft_pptx_continues_briefing_row_numbers_across_chunks(tmp_path):
    from pptx import Presentation

    rows = [REQUIRED_HEADER.strip()]
    image_dir = tmp_path / "images"
    image_dir.mkdir()
    photo_batches: list[DraftPhotoBatch] = []
    for idx in range(1, 7):
        offer_id = f"offer_{idx:02d}"
        rows.append(_offer_csv_row(offer_id))
        photo_batches.append(
            DraftPhotoBatch(
                tag=offer_id,
                image_paths=[_write_image(image_dir / f"briefing-{offer_id}.png")],
            )
        )
    csv_path = _write_csv(tmp_path / "offers.csv", "\n".join(rows) + "\n")
    payload = build_draft_payload(csv_path, photo_batches=photo_batches)

    output_path = tmp_path / "draft-briefing-numbering.pptx"
    create_draft_pptx(payload, output_path)

    prs = Presentation(str(output_path))
    second_briefing_slide = prs.slides[3]
    briefing_table = next(shape.table for shape in second_briefing_slide.shapes if getattr(shape, "has_table", False))

    assert briefing_table.cell(1, 0).text == "6"



def test_create_draft_pptx_writes_expected_slides_for_basic_payload(tmp_path):
    csv_path = _write_csv(
        tmp_path / "offers.csv",
        REQUIRED_HEADER
        + "offer_01,캐슬프라자 3층,성수동 2가 289-13 / 성수역 도보 2분,3층 / 18평,보증금 2000 / 월세 160 / 관리비 42,주차 1대; 인테리어 있음; 공실,offer_01\n"
        + "offer_02,동성빌딩 4층,성수동 2가 277-50 / 성수역 도보 4분,4층 / 14평,보증금 1100 / 월세 105 / 관리비 16,주차 가능; 공실; 화물 EV 가능,offer_02\n",
    )
    image_dir = tmp_path / "images"
    image_dir.mkdir()

    payload = build_draft_payload(
        csv_path,
        photo_batches=[
            DraftPhotoBatch(
                tag="offer_01",
                image_paths=[
                    _write_image(image_dir / "offer_01_01.png"),
                    _write_image(image_dir / "offer_01_02.png"),
                    _write_image(image_dir / "offer_01_03.png"),
                ],
            ),
            DraftPhotoBatch(
                tag="offer_02",
                image_paths=[
                    _write_image(image_dir / "offer_02_01.png"),
                    _write_image(image_dir / "offer_02_02.png"),
                    _write_image(image_dir / "offer_02_03.png"),
                    _write_image(image_dir / "offer_02_04.png"),
                    _write_image(image_dir / "offer_02_05.png"),
                ],
            ),
        ],
    )

    output_path = tmp_path / "draft-basic.pptx"
    created = create_draft_pptx(payload, output_path, title="성수 임대차 제안서", client="OO브랜드")

    assert Path(created) == output_path
    assert output_path.exists()
    assert _count_slides(output_path) == 9



def test_create_draft_pptx_uses_fixed_intro_pages_and_pdf_like_post_intro_structure(tmp_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Inches

    csv_path = _write_csv(
        tmp_path / "offers.csv",
        REQUIRED_HEADER
        + "offer_01,서울숲드림타워,성수동2가323,2층 / 12평,보증금 2200 / 월세 220 / 관리비 35,성수역 도보 6분; 무료주차 1대 가능; 지상 로딩도크 및 화물 E/V 有,offer_01\n",
    )
    image_dir = tmp_path / "images"
    image_dir.mkdir()

    payload = build_draft_payload(
        csv_path,
        photo_batches=[
            DraftPhotoBatch(
                tag="offer_01",
                image_paths=[
                    _write_image(image_dir / "offer_01_main.png", size=(1600, 900), color=(40, 70, 120)),
                    _write_image(image_dir / "offer_01_sub.png", size=(1200, 800), color=(180, 190, 205)),
                    _write_image(image_dir / "offer_01_map.png", size=(900, 1400), color=(230, 230, 230)),
                ],
            ),
        ],
    )

    output_path = tmp_path / "draft-single-offer.pptx"
    create_draft_pptx(payload, output_path, title="성수 임대차 제안서", client="OO브랜드")

    prs = Presentation(str(output_path))
    intro_slide = prs.slides[0]
    profile_slide = prs.slides[1]
    briefing_slide = prs.slides[2]
    overview_slide = prs.slides[3]
    detail_slide = prs.slides[4]

    intro_shape_types = [shape.shape_type for shape in intro_slide.shapes]
    profile_shape_types = [shape.shape_type for shape in profile_slide.shapes]
    intro_text = "\n".join(
        paragraph.text
        for shape in intro_slide.shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
        if paragraph.text
    )
    profile_text = "\n".join(
        paragraph.text
        for shape in profile_slide.shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
        if paragraph.text
    )
    briefing_text = "\n".join(
        paragraph.text
        for shape in briefing_slide.shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
        if paragraph.text
    )
    overview_text = "\n".join(
        paragraph.text
        for shape in overview_slide.shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
        if paragraph.text
    )
    detail_text = "\n".join(
        paragraph.text
        for shape in detail_slide.shapes
        if getattr(shape, "has_text_frame", False)
        for paragraph in shape.text_frame.paragraphs
        if paragraph.text
    )

    assert len(prs.slides) == 6
    assert intro_shape_types.count(MSO_SHAPE_TYPE.PICTURE) >= 1
    assert profile_shape_types.count(MSO_SHAPE_TYPE.PICTURE) >= 1
    assert intro_text == ""
    assert profile_text == ""
    assert "01. 추천 매물 브리핑" in briefing_text
    assert any(getattr(shape, "has_table", False) for shape in briefing_slide.shapes)
    assert "02. 매물 소개 (1)" in overview_text
    assert "보증금" in overview_text
    assert "서울숲드림타워" in detail_text
    assert "보증금" not in detail_text
    detail_pictures = [
        shape
        for shape in detail_slide.shapes
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE and shape.top > Inches(1.5)
    ]
    assert len(detail_pictures) == 2



def test_create_draft_pptx_uses_brand_green_system_for_offer_and_closing_slides(tmp_path):
    from pptx import Presentation
    from pptx.util import Inches

    csv_path = _write_csv(
        tmp_path / "offers.csv",
        REQUIRED_HEADER
        + "offer_01,서울숲드림타워,성수동2가323,2층 / 12평,보증금 2200 / 월세 220 / 관리비 35,성수역 도보 6분; 무료주차 1대 가능; 지상 로딩도크 및 화물 E/V 有,offer_01\n",
    )
    image_dir = tmp_path / "images"
    image_dir.mkdir()

    payload = build_draft_payload(
        csv_path,
        photo_batches=[
            DraftPhotoBatch(
                tag="offer_01",
                image_paths=[
                    _write_image(image_dir / "offer_01_main.png", size=(1600, 900), color=(40, 70, 120)),
                    _write_image(image_dir / "offer_01_sub.png", size=(1200, 800), color=(180, 190, 205)),
                    _write_image(image_dir / "offer_01_map.png", size=(900, 1400), color=(230, 230, 230)),
                ],
            ),
        ],
    )

    output_path = tmp_path / "draft-brand-green.pptx"
    create_draft_pptx(payload, output_path)

    prs = Presentation(str(output_path))
    briefing_slide = prs.slides[2]
    offer_slide = prs.slides[3]
    gallery_slide = prs.slides[4]
    closing_slide = prs.slides[5]

    offer_header_band = next(
        shape for shape in offer_slide.shapes if shape.width == prs.slide_width and shape.height > 600000
    )
    gallery_header_band = next(
        shape for shape in gallery_slide.shapes if shape.width == prs.slide_width and shape.height > 600000
    )
    closing_background = next(
        shape for shape in closing_slide.shapes if shape.width == prs.slide_width and shape.height == prs.slide_height
    )
    briefing_header_band = next(
        shape for shape in briefing_slide.shapes if shape.width == prs.slide_width and shape.height > 600000
    )
    offer_header_tab = next(
        shape
        for shape in offer_slide.shapes
        if shape.left == Inches(0.554)
        and shape.top == Inches(0.663)
        and shape.width == Inches(3.325)
        and shape.height == Inches(0.554)
    )

    assert offer_header_band.fill.fore_color.rgb == gallery_header_band.fill.fore_color.rgb == closing_background.fill.fore_color.rgb
    assert briefing_header_band.fill.fore_color.rgb == offer_header_band.fill.fore_color.rgb
    assert str(offer_header_band.fill.fore_color.rgb) == "1D4427"
    assert str(offer_header_tab.fill.fore_color.rgb) == "FFFFFF"



def test_create_draft_pptx_repeats_sample_pptx_header_asset_across_content_slides(tmp_path):
    from pptx import Presentation
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_SHAPE_TYPE
    from pptx.util import Inches, Pt

    header_logo_asset = Path(__file__).resolve().parents[2] / "assets" / "ppt_draft" / "linchpin_header_logo.emf"
    assert header_logo_asset.exists()

    csv_path = _write_csv(
        tmp_path / "offers.csv",
        REQUIRED_HEADER
        + "offer_01,서울숲드림타워,성수동2가323,2층 / 12평,보증금 2200 / 월세 220 / 관리비 35,성수역 도보 6분; 무료주차 1대 가능; 지상 로딩도크 및 화물 E/V 有,offer_01\n",
    )
    image_dir = tmp_path / "images"
    image_dir.mkdir()

    payload = build_draft_payload(
        csv_path,
        photo_batches=[
            DraftPhotoBatch(
                tag="offer_01",
                image_paths=[_write_image(image_dir / f"{idx:02d}.png", size=(1600, 900), color=(40 + idx, 70, 120)) for idx in range(8)],
            ),
        ],
    )

    output_path = tmp_path / "draft-sample-header.pptx"
    create_draft_pptx(payload, output_path)

    prs = Presentation(str(output_path))
    content_slides = [prs.slides[2], prs.slides[3], prs.slides[4], prs.slides[5]]

    for slide in content_slides:
        header_band = next(shape for shape in slide.shapes if shape.width == prs.slide_width and shape.top == 0)
        assert header_band.left == 0
        assert header_band.height == Inches(1.217)
        assert str(header_band.fill.fore_color.rgb) == "1D4427"

        tab = next(
            shape
            for shape in slide.shapes
            if getattr(shape, "auto_shape_type", None) == MSO_AUTO_SHAPE_TYPE.ROUND_1_RECTANGLE
            and shape.left == Inches(0.554)
            and shape.top == Inches(0.663)
        )
        assert tab.width == Inches(3.325)
        assert tab.height == Inches(0.554)
        assert str(tab.fill.fore_color.rgb) == "FFFFFF"

        title_box = next(
            shape
            for shape in slide.shapes
            if getattr(shape, "has_text_frame", False)
            and any(paragraph.text.startswith(("01.", "02.", "03.")) for paragraph in shape.text_frame.paragraphs)
        )
        assert title_box.left == Inches(0.625)
        assert title_box.top == Inches(0.764)
        assert title_box.height == Inches(0.438)
        title_paragraph = title_box.text_frame.paragraphs[0]
        assert title_paragraph.font.name == "-윤고딕310"
        assert title_paragraph.font.size == Pt(20)
        assert title_paragraph.font.bold is True

        tab_line = next(
            shape
            for shape in slide.shapes
            if shape.left == Inches(0.723)
            and shape.top == Inches(1.217)
            and shape.width == Inches(2.964)
            and shape.height == 0
        )
        assert str(tab_line.line.color.rgb) == "1D4427"

        logo = next(
            shape
            for shape in slide.shapes
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE
            and shape.left == Inches(10.283)
            and shape.top == Inches(0.425)
        )
        assert logo.width == Inches(2.644)
        assert logo.height == Inches(0.455)

        assert not any(
            getattr(shape, "has_text_frame", False)
            and any(paragraph.text == "LINCHPIN" for paragraph in shape.text_frame.paragraphs)
            for shape in slide.shapes
        )



def test_create_draft_pptx_splits_large_gallery_across_multiple_slides(tmp_path):
    csv_path = _write_csv(
        tmp_path / "offers.csv",
        REQUIRED_HEADER
        + "offer_01,캐슬프라자 7층,성수동 2가 289-13 / 성수역 도보 2분,7층 / 21평,보증금 2000 / 월세 160 / 관리비 42,주차 1대; 인테리어 있음; 공실,offer_01\n",
    )
    image_dir = tmp_path / "images"
    image_dir.mkdir()

    payload = build_draft_payload(
        csv_path,
        photo_batches=[
            DraftPhotoBatch(
                tag="offer_01",
                image_paths=[_write_image(image_dir / f"{idx:02d}.png") for idx in range(10)],
            ),
        ],
    )

    output_path = tmp_path / "draft-gallery-split.pptx"
    create_draft_pptx(payload, output_path)

    assert output_path.exists()
    assert _count_slides(output_path) == 7
