"""Spreadsheet parsing, photo-tag matching, and PPT draft generation helpers."""
from __future__ import annotations

import csv
import math
import posixpath
import re
import xml.etree.ElementTree as ET
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Iterable
from zipfile import BadZipFile, ZipFile

from gateway.ppt_draft_state import DraftPhotoBatch


REQUIRED_COLUMNS = (
    "id",
    "name",
    "location",
    "size_floor",
    "price",
    "points",
    "photo_tag",
)
_XLSX_NS = {
    "main": "http://schemas.openxmlformats.org/spreadsheetml/2006/main",
    "rel": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}
_INTRO_ASSET_DIR = Path(__file__).resolve().parents[1] / "assets" / "ppt_draft"
_FIXED_INTRO_PAGE_1 = _INTRO_ASSET_DIR / "lynchpin_intro_page1_unified.png"
_FIXED_INTRO_PAGE_2 = _INTRO_ASSET_DIR / "lynchpin_intro_page2_unified.png"
_HEADER_LOGO_ASSET = _INTRO_ASSET_DIR / "linchpin_header_logo.png"


class DraftInputError(Exception):
    def __init__(self, code: str, message: str, *, details: dict | None = None):
        super().__init__(message)
        self.code = code
        self.details = details or {}


@dataclass(slots=True)
class ParsedOffer:
    id: str
    name: str
    location: str
    size_floor: str
    price: str
    points: list[str]
    photo_tag: str
    photo_paths: list[str] = field(default_factory=list)


@dataclass(slots=True)
class DraftBuildPayload:
    csv_path: str
    offers: list[ParsedOffer]
    matched_photo_count: int
    unmatched_photo_tags: list[str]


@dataclass(slots=True)
class DraftTemplateConfig:
    briefing_offers_per_slide: int = 5
    detail_images_per_slide: int = 2
    gallery_images_per_slide: int = 9
    intro_assets: tuple[Path, Path] = (_FIXED_INTRO_PAGE_1, _FIXED_INTRO_PAGE_2)

    def __post_init__(self) -> None:
        try:
            normalized_intro_assets = tuple(Path(path) for path in self.intro_assets)
        except TypeError as exc:
            raise ValueError("intro_assets must be path-like values") from exc
        object.__setattr__(self, "intro_assets", normalized_intro_assets)

        if self.briefing_offers_per_slide <= 0:
            raise ValueError("briefing_offers_per_slide must be positive")
        if self.gallery_images_per_slide <= 0:
            raise ValueError("gallery_images_per_slide must be positive")
        if self.detail_images_per_slide != 2:
            raise ValueError("detail_images_per_slide currently supports exactly 2 images per slide")
        if len(self.intro_assets) != 2:
            raise ValueError("intro_assets must contain exactly two fixed intro asset paths")


@dataclass(slots=True)
class OfferRenderPlan:
    offer: ParsedOffer
    hero_image: str | None
    map_image: str | None
    detail_images: list[str]
    gallery_chunks: list[list[str]]


@dataclass(slots=True)
class DeckRenderPlan:
    briefing_chunks: list[list[ParsedOffer]]
    offer_plans: list[OfferRenderPlan]


@dataclass(slots=True)
class OfferOverviewLayout:
    headline_text: str
    title_left: float
    title_top: float
    title_width: float
    title_height: float
    title_font_size: int
    subtitle_left: float
    subtitle_top: float
    subtitle_width: float
    price_row_top: float
    price_label_lefts: tuple[float, float, float]
    note_label_left: float
    note_value_left: float
    note_row_top: float
    photos_top: float
    photo_height: float


def _normalize_tag(value: str) -> str:
    return str(value or "").strip().lower()


def _split_points(value: str) -> list[str]:
    normalized = str(value or "").replace("\r\n", "\n").replace("\r", "\n")
    parts: list[str] = []
    for line in normalized.split("\n"):
        for chunk in line.split(";"):
            cleaned = chunk.strip()
            if cleaned:
                parts.append(cleaned)
    return parts


def _column_index_from_cell_reference(reference: str) -> int | None:
    letters = "".join(ch for ch in str(reference or "") if ch.isalpha())
    if not letters:
        return None
    index = 0
    for ch in letters.upper():
        index = index * 26 + (ord(ch) - 64)
    return index - 1


def _read_xlsx_shared_strings(zf: ZipFile) -> list[str]:
    if "xl/sharedStrings.xml" not in zf.namelist():
        return []
    root = ET.fromstring(zf.read("xl/sharedStrings.xml"))
    values: list[str] = []
    for item in root.findall("main:si", _XLSX_NS):
        values.append("".join(node.text or "" for node in item.findall(".//main:t", _XLSX_NS)))
    return values


def _xlsx_cell_value(cell: ET.Element, shared_strings: list[str]) -> str:
    cell_type = cell.attrib.get("t", "")
    if cell_type == "inlineStr":
        return "".join(node.text or "" for node in cell.findall(".//main:t", _XLSX_NS))

    value_node = cell.find("main:v", _XLSX_NS)
    value = value_node.text if value_node is not None and value_node.text is not None else ""
    if cell_type == "s":
        try:
            index = int(value)
        except (TypeError, ValueError):
            return ""
        return shared_strings[index] if 0 <= index < len(shared_strings) else ""
    return value


def _read_xlsx_rows(path: Path) -> tuple[list[str], list[dict[str, str]]]:
    try:
        with ZipFile(path) as zf:
            names = set(zf.namelist())
            if "xl/workbook.xml" not in names:
                raise DraftInputError("invalid_xlsx", f"Workbook is missing xl/workbook.xml: {path.name}")
            if "xl/_rels/workbook.xml.rels" not in names:
                raise DraftInputError("invalid_xlsx", f"Workbook is missing xl/_rels/workbook.xml.rels: {path.name}")

            workbook_root = ET.fromstring(zf.read("xl/workbook.xml"))
            rels_root = ET.fromstring(zf.read("xl/_rels/workbook.xml.rels"))
            rel_map: dict[str, str] = {}
            for rel in rels_root:
                rel_id = rel.attrib.get("Id")
                target = rel.attrib.get("Target")
                if not rel_id or not target:
                    continue
                if target.startswith("/"):
                    resolved = target.lstrip("/")
                else:
                    resolved = posixpath.normpath(posixpath.join("xl", target))
                rel_map[rel_id] = resolved

            sheets = workbook_root.findall("main:sheets/main:sheet", _XLSX_NS)
            if not sheets:
                return [], []
            selected_sheet = next(
                (sheet for sheet in sheets if str(sheet.attrib.get("name") or "").strip().lower() == "in"),
                sheets[0],
            )
            rel_id = selected_sheet.attrib.get("{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id")
            sheet_path = rel_map.get(rel_id or "")
            if not sheet_path or sheet_path not in names:
                raise DraftInputError("invalid_xlsx", f"Worksheet XML could not be resolved for: {path.name}")

            worksheet_root = ET.fromstring(zf.read(sheet_path))
            shared_strings = _read_xlsx_shared_strings(zf)
            rows: list[list[str]] = []
            for row in worksheet_root.findall("main:sheetData/main:row", _XLSX_NS):
                indexed_values: dict[int, str] = {}
                max_index = -1
                for fallback_index, cell in enumerate(row.findall("main:c", _XLSX_NS)):
                    column_index = _column_index_from_cell_reference(cell.attrib.get("r", ""))
                    if column_index is None:
                        column_index = fallback_index
                    indexed_values[column_index] = _xlsx_cell_value(cell, shared_strings)
                    max_index = max(max_index, column_index)
                if max_index >= 0:
                    rows.append([indexed_values.get(index, "") for index in range(max_index + 1)])
    except BadZipFile as exc:
        raise DraftInputError("invalid_xlsx", f"Offers workbook is not a valid .xlsx file: {path.name}") from exc
    except ET.ParseError as exc:
        raise DraftInputError("invalid_xlsx", f"Offers workbook XML is malformed: {path.name}") from exc

    if not rows:
        return [], []

    fieldnames = [str(value or "").strip() for value in rows[0]]
    dict_rows: list[dict[str, str]] = []
    for values in rows[1:]:
        dict_rows.append(
            {
                fieldnames[index]: str(values[index] or "") if index < len(values) else ""
                for index in range(len(fieldnames))
                if fieldnames[index]
            }
        )
    return fieldnames, dict_rows


def _read_offer_rows(path: Path) -> tuple[list[str], list[dict[str, str]]]:
    suffix = path.suffix.lower()
    if suffix == ".xlsx":
        return _read_xlsx_rows(path)
    if suffix != ".csv":
        raise DraftInputError(
            "unsupported_offers_file",
            f"Offers file must be a .csv or .xlsx workbook: {path.name}",
            details={"path": str(path)},
        )

    with path.open("r", encoding="utf-8-sig", newline="") as handle:
        reader = csv.DictReader(handle)
        fieldnames = [str(name or "").strip() for name in (reader.fieldnames or [])]
        rows: list[dict[str, str]] = []
        for row in reader:
            normalized_row: dict[str, str] = {}
            for key, value in row.items():
                normalized_key = str(key or "").strip()
                if not normalized_key:
                    continue
                normalized_row[normalized_key] = str(value or "")
            rows.append(normalized_row)
    return fieldnames, rows



def parse_offers_csv(csv_path: str | Path) -> list[ParsedOffer]:
    path = Path(csv_path)
    fieldnames, rows = _read_offer_rows(path)
    missing = [column for column in REQUIRED_COLUMNS if column not in fieldnames]
    if missing:
        raise DraftInputError(
            "missing_columns",
            f"Missing required columns: {', '.join(missing)}",
            details={"missing_columns": missing},
        )

    offers: list[ParsedOffer] = []
    for row in rows:
        if not row:
            continue
        if not any(str(value or "").strip() for value in row.values()):
            continue
        offers.append(
            ParsedOffer(
                id=str(row.get("id") or "").strip(),
                name=str(row.get("name") or "").strip(),
                location=str(row.get("location") or "").strip(),
                size_floor=str(row.get("size_floor") or "").strip(),
                price=str(row.get("price") or "").strip(),
                points=_split_points(str(row.get("points") or "")),
                photo_tag=str(row.get("photo_tag") or "").strip(),
            )
        )
    return offers



def _aggregate_photo_batches(photo_batches: Iterable[DraftPhotoBatch]) -> dict[str, list[str]]:
    tag_to_paths: dict[str, list[str]] = {}
    for batch in photo_batches:
        tag = _normalize_tag(batch.tag)
        if not tag:
            continue
        paths = tag_to_paths.setdefault(tag, [])
        for path in batch.image_paths:
            value = str(path).strip()
            if value:
                paths.append(value)
    return tag_to_paths



def build_draft_payload(
    csv_path: str | Path,
    *,
    photo_batches: Iterable[DraftPhotoBatch],
) -> DraftBuildPayload:
    offers = parse_offers_csv(csv_path)
    tag_to_paths = _aggregate_photo_batches(photo_batches)
    required_tags = {_normalize_tag(offer.photo_tag) for offer in offers}

    unmatched_required = [
        offer.photo_tag
        for offer in offers
        if _normalize_tag(offer.photo_tag) not in tag_to_paths
    ]
    if unmatched_required:
        raise DraftInputError(
            "photo_tag_not_matched",
            f"No photo batch matched required tag(s): {', '.join(unmatched_required)}",
            details={"unmatched_required_tags": unmatched_required},
        )

    for offer in offers:
        offer.photo_paths = list(tag_to_paths.get(_normalize_tag(offer.photo_tag), []))

    unmatched_photo_tags = sorted(
        tag for tag in tag_to_paths.keys()
        if tag not in required_tags
    )
    matched_photo_count = sum(len(offer.photo_paths) for offer in offers)

    return DraftBuildPayload(
        csv_path=str(Path(csv_path)),
        offers=offers,
        matched_photo_count=matched_photo_count,
        unmatched_photo_tags=unmatched_photo_tags,
    )



def _safe_image_size(image_path: str | Path) -> tuple[int, int] | None:
    try:
        from PIL import Image

        with Image.open(str(image_path)) as img:
            return img.size
    except Exception:
        return None



def choose_overview_images(image_paths: list[str]) -> tuple[str | None, str | None]:
    normalized = [str(path).strip() for path in image_paths if str(path).strip()]
    if not normalized:
        return None, None

    portraits: list[str] = []
    others: list[str] = []
    for image_path in normalized:
        size = _safe_image_size(image_path)
        if size and size[0] < size[1]:
            portraits.append(image_path)
        else:
            others.append(image_path)

    hero_path = others[0] if others else normalized[0]
    secondary_path = portraits[0] if portraits else next((path for path in normalized if path != hero_path), None)
    if secondary_path == hero_path:
        secondary_path = next((path for path in normalized if path != hero_path), None)
    return hero_path, secondary_path



def choose_detail_images(
    image_paths: list[str],
    *,
    excluded_paths: list[str] | None = None,
    count: int = 2,
) -> list[str]:
    normalized: list[str] = []
    for image_path in image_paths:
        value = str(image_path).strip()
        if value and value not in normalized:
            normalized.append(value)
    if not normalized or count <= 0:
        return []

    excluded = {str(path).strip() for path in (excluded_paths or []) if str(path).strip()}
    unique_candidates = [path for path in normalized if path not in excluded]

    preferred: list[str] = []
    fallback: list[str] = []
    for image_path in unique_candidates:
        size = _safe_image_size(image_path)
        if size and size[0] >= size[1]:
            preferred.append(image_path)
        else:
            fallback.append(image_path)

    selected: list[str] = []
    for image_path in preferred + fallback:
        if image_path not in selected:
            selected.append(image_path)
        if len(selected) == count:
            return selected

    for image_path in normalized:
        if image_path not in selected:
            selected.append(image_path)
        if len(selected) == count:
            break
    return selected



def remaining_gallery_images(image_paths: list[str], *, excluded_paths: list[str] | None = None) -> list[str]:
    excluded = {str(path).strip() for path in (excluded_paths or []) if str(path).strip()}
    remaining: list[str] = []
    for image_path in image_paths:
        value = str(image_path).strip()
        if not value or value in excluded or value in remaining:
            continue
        remaining.append(value)
    return remaining



def _chunk(items: list, size: int) -> list[list]:
    return [items[index:index + size] for index in range(0, len(items), size)]



def build_offer_render_plan(offer: ParsedOffer, config: DraftTemplateConfig) -> OfferRenderPlan:
    hero_image, map_image = choose_overview_images(offer.photo_paths)
    detail_images = choose_detail_images(
        offer.photo_paths,
        excluded_paths=[hero_image or "", map_image or ""],
        count=config.detail_images_per_slide,
    )
    gallery_images = remaining_gallery_images(
        offer.photo_paths,
        excluded_paths=[hero_image or "", map_image or "", *detail_images],
    )
    return OfferRenderPlan(
        offer=offer,
        hero_image=hero_image,
        map_image=map_image,
        detail_images=detail_images,
        gallery_chunks=_chunk(gallery_images, config.gallery_images_per_slide),
    )



def build_deck_render_plan(payload: DraftBuildPayload, config: DraftTemplateConfig) -> DeckRenderPlan:
    return DeckRenderPlan(
        briefing_chunks=_chunk(payload.offers, config.briefing_offers_per_slide),
        offer_plans=[build_offer_render_plan(offer, config) for offer in payload.offers],
    )



def compose_offer_headline(offer: ParsedOffer) -> str:
    return f"{offer.location} {offer.name}".strip()



def _wrap_offer_headline(text: str, *, target_chars: int = 30, min_split: int = 18, max_split: int = 40) -> str:
    normalized = " ".join(str(text or "").split())
    if len(normalized) <= target_chars:
        return normalized

    separators = [" / ", " | ", " - ", " "]
    candidates: list[int] = []
    for separator in separators:
        start = 0
        while True:
            idx = normalized.find(separator, start)
            if idx == -1:
                break
            split_at = idx + len(separator.strip())
            if min_split <= split_at <= max_split:
                candidates.append(split_at)
            start = idx + len(separator)

    if candidates:
        split_at = min(candidates, key=lambda idx: abs(idx - target_chars))
    else:
        split_at = min(len(normalized), target_chars)

    left = normalized[:split_at].rstrip(" /|-")
    right = normalized[split_at:].lstrip(" /|-")
    if not right:
        return normalized
    return f"{left}\n{right}"



def format_price_value(value: str) -> str:
    """Add thousands separators to integer money amounts while preserving units."""
    text = str(value or "").strip()

    def replace_number(match: re.Match[str]) -> str:
        raw = match.group(0)
        digits = raw.replace(",", "")
        if not digits.isdigit():
            return raw
        number = int(digits)
        if number < 1000:
            return raw
        return f"{number:,}"

    return re.sub(r"(?<![\d.])\d[\d,]*(?![\d.])", replace_number, text)



def build_offer_overview_layout(offer: ParsedOffer) -> OfferOverviewLayout:
    headline = _wrap_offer_headline(compose_offer_headline(offer), target_chars=28, min_split=16, max_split=36)
    is_wrapped = "\n" in headline
    return OfferOverviewLayout(
        headline_text=headline,
        title_left=0.894,
        title_top=1.58 if is_wrapped else 1.792,
        title_width=6.35,
        title_height=0.78 if is_wrapped else 0.404,
        title_font_size=20 if is_wrapped else 22,
        subtitle_left=0.894,
        subtitle_top=2.34 if is_wrapped else 2.128,
        subtitle_width=5.042,
        price_row_top=1.673,
        price_label_lefts=(7.473, 9.129, 10.579),
        note_label_left=7.473,
        note_value_left=8.406,
        note_row_top=2.107,
        photos_top=2.705,
        photo_height=4.795,
    )



def create_draft_pptx(
    payload: DraftBuildPayload,
    output_path: str | Path,
    *,
    title: str | None = None,
    client: str | None = None,
    template_config: DraftTemplateConfig | None = None,
) -> str:
    import re

    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt

    FONT_FAMILY = "맑은 고딕"
    HEADER_FONT_FAMILY = "-윤고딕310"
    COLOR_PRIMARY = "1D4427"
    COLOR_PRIMARY_DARK = "1D4427"
    COLOR_ACCENT = "1D4427"
    COLOR_TEXT = "24312C"
    COLOR_LIGHT = "F4F6F3"
    COLOR_BORDER = "D8DEDA"
    COLOR_WHITE = "FFFFFF"
    COLOR_MUTED = "6F7B75"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    config = template_config or DraftTemplateConfig()
    deck_plan = build_deck_render_plan(payload, config)

    def rgb(hex_code: str) -> RGBColor:
        return RGBColor(int(hex_code[0:2], 16), int(hex_code[2:4], 16), int(hex_code[4:6], 16))

    def style_paragraph(paragraph, *, size: int, color: str, bold: bool = False, alignment=None):
        paragraph.font.name = FONT_FAMILY
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = rgb(color)
        if alignment is not None:
            paragraph.alignment = alignment
        for run in paragraph.runs:
            run.font.name = FONT_FAMILY
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)

    def add_title(slide, text: str, *, top: float = 0.35, size: int = 24, color: str = COLOR_TEXT):
        box = slide.shapes.add_textbox(Inches(0.5), Inches(top), Inches(12.3), Inches(0.6))
        frame = box.text_frame
        frame.clear()
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        style_paragraph(paragraph, size=size, color=color, bold=True)
        return box

    def add_body_text(
        slide,
        text: str,
        *,
        left: float,
        top: float,
        width: float,
        height: float,
        size: int = 16,
        color: str = COLOR_TEXT,
    ):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        for idx, line in enumerate([part for part in text.split("\n") if part.strip()] or [""]):
            paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            paragraph.text = line
            style_paragraph(paragraph, size=size, color=color, bold=False)
        return box

    def add_points(slide, points: list[str], *, left: float, top: float, width: float):
        box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(2.2))
        frame = box.text_frame
        frame.clear()
        frame.word_wrap = True
        for idx, point in enumerate(points or ["포인트 없음"]):
            paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
            paragraph.text = point
            paragraph.level = 0
            paragraph.bullet = True
            style_paragraph(paragraph, size=15, color=COLOR_TEXT)
        return box

    def safe_image_size(image_path: str | Path) -> tuple[int, int] | None:
        return _safe_image_size(image_path)

    def add_picture_cover(slide, image_path: str | Path, *, left: float, top: float, width: float, height: float):
        picture = slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))
        size = safe_image_size(image_path)
        if not size or not size[0] or not size[1]:
            return picture
        image_ratio = size[0] / size[1]
        target_ratio = width / height
        if image_ratio > target_ratio:
            crop = (1 - (target_ratio / image_ratio)) / 2
            picture.crop_left = crop
            picture.crop_right = crop
        elif image_ratio < target_ratio:
            crop = (1 - (image_ratio / target_ratio)) / 2
            picture.crop_top = crop
            picture.crop_bottom = crop
        return picture

    def add_full_slide_image(slide, image_path: str | Path):
        return add_picture_cover(slide, image_path, left=0, top=0, width=13.333, height=7.5)

    def add_picture_grid(slide, image_paths: list[str], *, left: float, top: float, width: float, height: float):
        count = len(image_paths)
        if count <= 1:
            cols = 1
        elif count <= 4:
            cols = 2
        else:
            cols = 3
        rows = math.ceil(count / cols)
        gap_x = 0.12
        gap_y = 0.12
        cell_w = (width - gap_x * (cols - 1)) / cols
        cell_h = (height - gap_y * (rows - 1)) / rows
        for idx, image_path in enumerate(image_paths):
            row = idx // cols
            col = idx % cols
            x = left + col * (cell_w + gap_x)
            y = top + row * (cell_h + gap_y)
            add_picture_cover(slide, image_path, left=x, top=y, width=cell_w, height=cell_h)

    def add_label_tag(slide, text: str, *, left: float, top: float, width: float = 1.25):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(0.34),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(COLOR_PRIMARY_DARK)
        shape.line.color.rgb = rgb(COLOR_PRIMARY_DARK)
        frame = shape.text_frame
        frame.clear()
        frame.margin_left = Pt(6)
        frame.margin_right = Pt(6)
        frame.margin_top = Pt(2)
        frame.margin_bottom = Pt(1)
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        style_paragraph(paragraph, size=10, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
        return shape

    HEADER_BAND_HEIGHT = 1.217
    HEADER_TAB_LEFT = 0.554
    HEADER_TAB_TOP = 0.663
    HEADER_TAB_WIDTH = 3.325
    HEADER_TAB_HEIGHT = 0.554
    HEADER_TITLE_LEFT = 0.625
    HEADER_TITLE_TOP = 0.764
    HEADER_TITLE_WIDTH = 3.0
    HEADER_TITLE_HEIGHT = 0.438
    HEADER_UNDERLINE_LEFT = 0.723
    HEADER_UNDERLINE_TOP = 1.217
    HEADER_UNDERLINE_WIDTH = 2.964
    HEADER_LOGO_LEFT = 10.283
    HEADER_LOGO_TOP = 0.425
    HEADER_LOGO_WIDTH = 2.644
    HEADER_LOGO_HEIGHT = 0.455

    def style_header_paragraph(paragraph, *, size: int = 20, color: str = COLOR_TEXT, bold: bool = True):
        paragraph.font.name = HEADER_FONT_FAMILY
        paragraph.font.size = Pt(size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = rgb(color)
        for run in paragraph.runs:
            run.font.name = HEADER_FONT_FAMILY
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = rgb(color)

    def add_brand_lockup(slide):
        return slide.shapes.add_picture(
            str(_HEADER_LOGO_ASSET),
            Inches(HEADER_LOGO_LEFT),
            Inches(HEADER_LOGO_TOP),
            width=Inches(HEADER_LOGO_WIDTH),
            height=Inches(HEADER_LOGO_HEIGHT),
        )

    def add_small_dots(slide, *, left: float, top: float, count: int = 3):
        for index in range(count):
            dot = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                Inches(left + index * 0.13),
                Inches(top),
                Inches(0.07),
                Inches(0.07),
            )
            dot.fill.solid()
            dot.fill.fore_color.rgb = rgb(COLOR_PRIMARY)
            dot.line.color.rgb = rgb(COLOR_PRIMARY)

    def add_photo_placeholder(slide, text: str, *, left: float, top: float, width: float, height: float):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = rgb(COLOR_LIGHT)
        shape.line.color.rgb = rgb(COLOR_BORDER)
        frame = shape.text_frame
        frame.clear()
        frame.margin_left = Pt(14)
        frame.margin_right = Pt(14)
        frame.margin_top = Pt(14)
        frame.margin_bottom = Pt(14)
        paragraph = frame.paragraphs[0]
        paragraph.text = text
        style_paragraph(paragraph, size=18, color=COLOR_MUTED, bold=True, alignment=PP_ALIGN.CENTER)
        return shape

    def add_header_band(slide, section_label: str):
        band = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0),
            Inches(0),
            Inches(13.333),
            Inches(HEADER_BAND_HEIGHT),
        )
        band.fill.solid()
        band.fill.fore_color.rgb = rgb(COLOR_PRIMARY_DARK)
        band.line.color.rgb = rgb(COLOR_PRIMARY_DARK)

        tab = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUND_1_RECTANGLE,
            Inches(HEADER_TAB_LEFT),
            Inches(HEADER_TAB_TOP),
            Inches(HEADER_TAB_WIDTH),
            Inches(HEADER_TAB_HEIGHT),
        )
        tab.fill.solid()
        tab.fill.fore_color.rgb = rgb(COLOR_WHITE)
        tab.line.color.rgb = rgb(COLOR_WHITE)

        title_box = slide.shapes.add_textbox(
            Inches(HEADER_TITLE_LEFT),
            Inches(HEADER_TITLE_TOP),
            Inches(HEADER_TITLE_WIDTH),
            Inches(HEADER_TITLE_HEIGHT),
        )
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.margin_left = 0
        title_frame.margin_right = 0
        title_frame.margin_top = 0
        title_frame.margin_bottom = 0
        title_p = title_frame.paragraphs[0]
        title_p.text = section_label
        style_header_paragraph(title_p)

        underline = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            Inches(HEADER_UNDERLINE_LEFT),
            Inches(HEADER_UNDERLINE_TOP),
            Inches(HEADER_UNDERLINE_LEFT + HEADER_UNDERLINE_WIDTH),
            Inches(HEADER_UNDERLINE_TOP),
        )
        underline.line.color.rgb = rgb(COLOR_PRIMARY_DARK)
        underline.line.width = Pt(0.75)

        add_brand_lockup(slide)
        return band

    def parse_price_components(price_text: str) -> list[tuple[str, str]]:
        alias_map = {
            "보증금": "보증금",
            "월세": "월임차료",
            "월임차료": "월임차료",
            "월임대료": "월임차료",
            "관리비": "관리비",
        }
        results: list[tuple[str, str]] = []
        for raw_segment in re.split(r"[\/\n]", str(price_text or "")):
            segment = raw_segment.strip()
            if not segment:
                continue
            label = next((canonical for key, canonical in alias_map.items() if segment.startswith(key)), None)
            if label is None:
                continue
            value = segment
            for key, canonical in alias_map.items():
                if segment.startswith(key) and canonical == label:
                    value = segment[len(key):].strip(" :")
                    break
            if value:
                results.append((label, format_price_value(value)))
        if results:
            deduped: list[tuple[str, str]] = []
            seen: set[str] = set()
            for label, value in results:
                if label in seen:
                    continue
                seen.add(label)
                deduped.append((label, value))
            return deduped
        return [("가격", format_price_value(str(price_text or "-").strip() or "-"))]

    def price_component_map(price_text: str) -> dict[str, str]:
        return {label: value for label, value in parse_price_components(price_text)}

    def format_size_floor(size_floor: str) -> str:
        parts = [part.strip() for part in str(size_floor or "").split("/") if part.strip()]
        if len(parts) >= 2:
            return f"{parts[0]}   |   전용 약 {parts[1]}"
        return str(size_floor or "-")

    def format_area_for_table(size_floor: str) -> str:
        parts = [part.strip() for part in str(size_floor or "").split("/") if part.strip()]
        if len(parts) >= 2:
            return parts[1].replace("평", "py").replace(" ", "")
        return str(size_floor or "-")

    def add_label_value_pair(
        slide,
        label: str,
        value: str,
        *,
        left: float,
        top: float,
        label_width: float | None = None,
        value_width: float | None = None,
    ) -> float:
        computed_label_width = label_width or max(0.88, min(1.25, 0.46 + len(label) * 0.15))
        computed_value_width = value_width or max(0.8, min(1.75, 0.38 + len(value) * 0.13))
        pill = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(left),
            Inches(top),
            Inches(computed_label_width),
            Inches(0.34),
        )
        pill.fill.solid()
        pill.fill.fore_color.rgb = rgb(COLOR_PRIMARY_DARK)
        pill.line.color.rgb = rgb(COLOR_PRIMARY_DARK)
        frame = pill.text_frame
        frame.clear()
        frame.margin_left = Pt(8)
        frame.margin_right = Pt(8)
        label_p = frame.paragraphs[0]
        label_p.text = label
        style_paragraph(label_p, size=10, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        value_box = slide.shapes.add_textbox(
            Inches(left + computed_label_width + 0.08),
            Inches(top - 0.02),
            Inches(computed_value_width),
            Inches(0.38),
        )
        value_frame = value_box.text_frame
        value_frame.clear()
        value_frame.margin_left = 0
        value_frame.margin_right = 0
        value_p = value_frame.paragraphs[0]
        value_p.text = value
        style_paragraph(value_p, size=15, color=COLOR_TEXT, bold=True)
        return left + computed_label_width + computed_value_width + 0.26

    def set_table_cell_text(cell, text: str, *, size: int, color: str, bold: bool = False, alignment=PP_ALIGN.CENTER):
        frame = cell.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.margin_left = Pt(5)
        frame.margin_right = Pt(5)
        frame.margin_top = Pt(5)
        frame.margin_bottom = Pt(5)
        lines = [line for line in str(text or "-").split("\n") if line.strip()] or ["-"]
        for index, line in enumerate(lines):
            paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
            paragraph.text = line
            style_paragraph(paragraph, size=size, color=color, bold=bold, alignment=alignment)

    def build_briefing_rows(offers: list[ParsedOffer], *, start_index: int = 1) -> list[list[str]]:
        rows: list[list[str]] = []
        for index, offer in enumerate(offers, start=start_index):
            prices = price_component_map(offer.price)
            remarks = offer.points[:3] if offer.points else ["-"]
            rows.append(
                [
                    str(index),
                    offer.location or "-",
                    offer.name or "-",
                    prices.get("보증금", "-"),
                    prices.get("월임차료", "-"),
                    prices.get("관리비", "-"),
                    format_area_for_table(offer.size_floor),
                    "\n".join(remarks),
                ]
            )
        return rows

    def add_briefing_slide(slide, offers: list[ParsedOffer], *, chunk_index: int = 1, row_start_index: int = 1):
        label = "01. 추천 매물 브리핑" if chunk_index == 1 else f"01. 추천 매물 브리핑 ({chunk_index})"
        add_header_band(slide, label)
        rows = build_briefing_rows(offers, start_index=row_start_index)
        headers = ["NO", "주소", "건물명", "보증금", "임차료", "관리비", "전용면적", "비고"]
        table_shape = slide.shapes.add_table(
            len(rows) + 1,
            len(headers),
            Inches(0.55),
            Inches(1.55),
            Inches(12.2),
            Inches(5.2),
        )
        table = table_shape.table
        col_widths = [0.55, 1.7, 1.65, 0.95, 0.95, 0.95, 1.05, 4.4]
        for index, width in enumerate(col_widths):
            table.columns[index].width = Inches(width)

        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.fill.solid()
            cell.fill.fore_color.rgb = rgb(COLOR_LIGHT)
            set_table_cell_text(cell, header, size=10, color=COLOR_PRIMARY_DARK, bold=True)
        table.rows[0].height = Inches(0.46)

        body_height = 4.74 / max(1, len(rows))
        for row_idx, row_values in enumerate(rows, start=1):
            table.rows[row_idx].height = Inches(min(0.92, max(0.64, body_height)))
            for col_idx, value in enumerate(row_values):
                cell = table.cell(row_idx, col_idx)
                cell.fill.solid()
                cell.fill.fore_color.rgb = rgb(COLOR_WHITE)
                align = PP_ALIGN.CENTER if col_idx != 7 else PP_ALIGN.LEFT
                size = 9 if col_idx == 7 else 10
                set_table_cell_text(cell, value, size=size, color=COLOR_TEXT, bold=False, alignment=align)
        return table_shape

    def add_offer_overview_slide(
        slide,
        offer: ParsedOffer,
        *,
        offer_index: int,
        representative_image: str | None,
        map_image: str | None,
    ):
        layout = build_offer_overview_layout(offer)

        add_header_band(slide, f"02. 매물 소개 ({offer_index})")
        add_small_dots(slide, left=0.92, top=1.50)

        name_box = slide.shapes.add_textbox(
            Inches(layout.title_left),
            Inches(layout.title_top),
            Inches(layout.title_width),
            Inches(layout.title_height),
        )
        name_frame = name_box.text_frame
        name_frame.clear()
        name_frame.word_wrap = True
        name_frame.margin_left = 0
        name_frame.margin_right = 0
        headline_lines = layout.headline_text.split("\n") if layout.headline_text else [""]
        for idx, line in enumerate(headline_lines):
            name_p = name_frame.paragraphs[0] if idx == 0 else name_frame.add_paragraph()
            name_p.text = line
            style_paragraph(name_p, size=layout.title_font_size, color=COLOR_TEXT, bold=True)

        subtitle_box = slide.shapes.add_textbox(
            Inches(layout.subtitle_left),
            Inches(layout.subtitle_top),
            Inches(layout.subtitle_width),
            Inches(0.34),
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.clear()
        subtitle_frame.margin_left = 0
        subtitle_frame.margin_right = 0
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.text = format_size_floor(offer.size_floor)
        style_paragraph(subtitle_p, size=15, color=COLOR_MUTED, bold=True)

        for left, (label, value) in zip(layout.price_label_lefts, parse_price_components(offer.price)[:3], strict=False):
            add_label_value_pair(
                slide,
                label,
                value,
                left=left,
                top=layout.price_row_top,
                label_width=0.855,
            )

        if offer.points:
            feature_pill = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(layout.note_label_left),
                Inches(layout.note_row_top),
                Inches(0.855),
                Inches(0.34),
            )
            feature_pill.fill.solid()
            feature_pill.fill.fore_color.rgb = rgb(COLOR_PRIMARY_DARK)
            feature_pill.line.color.rgb = rgb(COLOR_PRIMARY_DARK)
            feature_frame = feature_pill.text_frame
            feature_frame.clear()
            feature_frame.margin_left = Pt(6)
            feature_frame.margin_right = Pt(6)
            feature_label = feature_frame.paragraphs[0]
            feature_label.text = "비고"
            style_paragraph(feature_label, size=10, color=COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

            feature_box = slide.shapes.add_textbox(
                Inches(layout.note_value_left),
                Inches(layout.note_row_top - 0.02),
                Inches(4.35),
                Inches(0.5),
            )
            feature_text_frame = feature_box.text_frame
            feature_text_frame.clear()
            feature_text_frame.word_wrap = True
            feature_text_frame.margin_left = 0
            feature_text_frame.margin_right = 0
            feature_p = feature_text_frame.paragraphs[0]
            feature_p.text = " / ".join(offer.points[:3])
            style_paragraph(feature_p, size=11, color=COLOR_TEXT, bold=False)

        if representative_image:
            add_picture_cover(slide, representative_image, left=0.78, top=layout.photos_top, width=5.72, height=layout.photo_height)
        else:
            add_photo_placeholder(slide, "대표 사진 없음", left=0.78, top=layout.photos_top, width=5.72, height=layout.photo_height)

        if map_image:
            add_picture_cover(slide, map_image, left=6.84, top=layout.photos_top, width=5.72, height=layout.photo_height)
        else:
            secondary = next((path for path in offer.photo_paths if path != representative_image), None)
            if secondary:
                add_picture_cover(slide, secondary, left=6.84, top=layout.photos_top, width=5.72, height=layout.photo_height)
            else:
                add_photo_placeholder(slide, f"위치 정보\n{offer.location}", left=6.84, top=layout.photos_top, width=5.72, height=layout.photo_height)

    def add_offer_detail_slide(slide, offer: ParsedOffer, *, offer_index: int, detail_images: list[str]):
        add_header_band(slide, f"02. 매물 소개 ({offer_index})")
        add_small_dots(slide, left=0.92, top=1.38)

        name_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.58), Inches(8.0), Inches(0.5))
        name_frame = name_box.text_frame
        name_frame.clear()
        name_frame.margin_left = 0
        name_frame.margin_right = 0
        name_p = name_frame.paragraphs[0]
        name_p.text = f"{offer.location} {offer.name}".strip()
        style_paragraph(name_p, size=22, color=COLOR_TEXT, bold=True)

        if detail_images[:1]:
            add_picture_cover(slide, detail_images[0], left=0.82, top=2.25, width=5.7, height=4.45)
        else:
            add_photo_placeholder(slide, "상세 사진 없음", left=0.82, top=2.25, width=5.7, height=4.45)

        if len(detail_images) >= 2:
            add_picture_cover(slide, detail_images[1], left=6.82, top=2.25, width=5.7, height=4.45)
        else:
            add_photo_placeholder(slide, "추가 사진 없음", left=6.82, top=2.25, width=5.7, height=4.45)

    def add_offer_gallery_slide(slide, offer: ParsedOffer, *, offer_index: int, gallery_paths: list[str], chunk_index: int):
        label = f"03. 현장 사진 ({offer_index})" if chunk_index == 1 else f"03. 현장 사진 ({offer_index}-{chunk_index})"
        add_header_band(slide, label)
        add_small_dots(slide, left=0.92, top=1.38)
        title_box = slide.shapes.add_textbox(Inches(0.9), Inches(1.58), Inches(8.0), Inches(0.44))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.margin_left = 0
        title_frame.margin_right = 0
        title_p = title_frame.paragraphs[0]
        title_p.text = f"{offer.location} {offer.name}".strip()
        style_paragraph(title_p, size=20, color=COLOR_TEXT, bold=True)
        add_picture_grid(slide, gallery_paths, left=0.82, top=2.2, width=11.7, height=4.7)

    # Fixed intro pages from reference PDF
    intro_page_1, intro_page_2 = config.intro_assets
    if not intro_page_1.exists() or not intro_page_2.exists():
        missing = [str(path) for path in (intro_page_1, intro_page_2) if not path.exists()]
        raise DraftInputError(
            "missing_intro_assets",
            f"Fixed intro assets are missing: {', '.join(missing)}",
            details={"missing_assets": missing},
        )

    invalid_intro_assets = [
        str(path)
        for path in (intro_page_1, intro_page_2)
        if path.is_dir() or _safe_image_size(path) is None
    ]
    if invalid_intro_assets:
        raise DraftInputError(
            "invalid_intro_assets",
            f"Fixed intro assets must be readable image files: {', '.join(invalid_intro_assets)}",
            details={"invalid_assets": invalid_intro_assets},
        )

    slide = prs.slides.add_slide(blank)
    add_full_slide_image(slide, intro_page_1)

    slide = prs.slides.add_slide(blank)
    add_full_slide_image(slide, intro_page_2)

    # Briefing + offer detail sequence
    briefing_row_start = 1
    for chunk_index, offers_chunk in enumerate(deck_plan.briefing_chunks, start=1):
        slide = prs.slides.add_slide(blank)
        add_briefing_slide(
            slide,
            offers_chunk,
            chunk_index=chunk_index,
            row_start_index=briefing_row_start,
        )
        briefing_row_start += len(offers_chunk)

    for offer_index, offer_plan in enumerate(deck_plan.offer_plans, start=1):
        offer = offer_plan.offer

        slide = prs.slides.add_slide(blank)
        add_offer_overview_slide(
            slide,
            offer,
            offer_index=offer_index,
            representative_image=offer_plan.hero_image,
            map_image=offer_plan.map_image,
        )

        slide = prs.slides.add_slide(blank)
        add_offer_detail_slide(slide, offer, offer_index=offer_index, detail_images=offer_plan.detail_images)

        for chunk_index, gallery_chunk in enumerate(offer_plan.gallery_chunks, start=1):
            if not gallery_chunk:
                continue
            slide = prs.slides.add_slide(blank)
            add_offer_gallery_slide(slide, offer, offer_index=offer_index, gallery_paths=gallery_chunk, chunk_index=chunk_index)

    # Closing
    slide = prs.slides.add_slide(blank)
    closing_band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(7.5),
    )
    closing_band.fill.solid()
    closing_band.fill.fore_color.rgb = rgb(COLOR_PRIMARY_DARK)
    closing_band.line.color.rgb = rgb(COLOR_PRIMARY_DARK)
    add_title(slide, "검토 후 세부 조건은 PPT에서 수정해 사용하세요.", top=2.0, size=24, color=COLOR_WHITE)
    add_body_text(
        slide,
        "초안 자동 생성 완료\n- 가격/표현은 최종 검토\n- 사진 순서는 필요 시 수동 조정\n- 고객 맞춤 멘트만 추가하면 바로 사용 가능",
        left=1.0,
        top=3.0,
        width=11.0,
        height=2.0,
        size=18,
        color=COLOR_WHITE,
    )

    output = Path(output_path)
    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output))
    return str(output)
